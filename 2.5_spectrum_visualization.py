"""
@author: Alexander Studier-Fischer, Jan Odenthal, Berkin Oezdemir, University of Heidelberg
"""
###
# Displaying Spectrum and Histogram PNGs in Summary Powerpoint


from pptx import *
import os
from pathlib import Path
from pptx.util import Cm, Pt
import glob
from skimage.draw import line_aa
import pandas as pd
import numpy as np
from PIL import Image
from datetime import datetime

###

slides_per_pptx = 400
folder_in_use = '/_hypergui_1'

label = "_labelling_001"

use_label = True

start_date = ""
end_date = ""                       ### If new end date wanted, change here, make sure to use the same format (2018_01_01_00_00_00)



###

if start_date  != "" or end_date != "":
    tag = "_specified"
else:
    tag = ""
if start_date == "":
    start_date = "2018_01_01_00_00_01"
if end_date == "":
    end_date = datetime.now().strftime('%Y_%m_%d_%H_%M_%S')

all_paths = list()
all_labels = list()
if use_label == True:
    tag = "_specified"
    for i in glob.glob('./data/*/2*/' + label + '*.txt'):
        all_paths.append(os.path.abspath(Path(i).parent))
        all_labels.append(str(os.path.basename(i)[:-4]))
elif use_label == False:
    for i in glob.glob('./data/*/2*'):
        if os.path.basename(i) > start_date and os.path.basename(i) < end_date:
            all_paths.append(os.path.abspath(i))


print(all_labels)

home = os.getcwd().replace('\\', '/')

#from HyperGUI
def _draw_a_line(point1, point2, image):
    r0, c0 = point1
    r1, c1 = point2
    rr, cc, val = line_aa(c0, r0, c1, r1)
    for i in range(len(rr)):
        image[rr[i] % 480, cc[i] % 640] = (int(113 * val[i]), int(255 * val[i]), int(66 * val[i]))
    return image

def crop_coords(path):
    coords = np.asarray(pd.DataFrame(pd.read_csv(path + folder_in_use + '//MASK_COORDINATES.csv', delimiter= ',', header=None)), dtype=np.int)
    coords = np.append(coords, coords[0,:]).reshape(coords.shape[0]+1,2)
    crop_path_rgb= glob.glob(path +'/_crops/*_cropped_RGB-Image.png')[0]
    crop_path_ox = glob.glob(path + '/_crops/*_cropped_Oxygenation.png')[0]
    img_rgb = np.array(Image.open(crop_path_rgb))
    img_ox = np.array(Image.open(crop_path_ox))
    for i in range(len(coords)):
        if not i == len(coords)-1:
            _draw_a_line(coords[i,:], coords[i+1,:], img_rgb )
            _draw_a_line(coords[i, :], coords[i + 1, :], img_ox)
    pillow_image_rgb = Image.fromarray(img_rgb)
    pillow_image_ox = Image.fromarray(img_ox)
    pillow_image_rgb.save(os.path.abspath(Path(crop_path_rgb).parent) + '/' + os.path.basename(crop_path_rgb)[0:19] +'_RGB_cropped_with_mask.png')
    pillow_image_ox.save(os.path.abspath(Path(crop_path_ox).parent) + '/' + os.path.basename(crop_path_ox)[0:19] + '_Oxygenation_cropped_with_mask.png')

prs = Presentation()

MAIN_IMAGE_SIZE = (Cm(6.25), Cm(4.57))
width = MAIN_IMAGE_SIZE[0]
height = MAIN_IMAGE_SIZE[1]
UPPER_ONE = (Cm(0), Cm(3.15))
MIDDLE_ONE = (Cm(0), Cm(8.46))
BOTTOM_ONE = (Cm(0), Cm(13.77))
UPPER_TWO = (Cm(6.25), Cm(3.15))
MIDDLE_TWO = (Cm(6.25), Cm(8.46))
BOTTOM_TWO = (Cm(6.25), Cm(13.77))
UPPER_THREE = (Cm(12.5), Cm(3.15))
MIDDLE_THREE = (Cm(12.5), Cm(8.46))
BOTTOM_THREE = (Cm(12.5), Cm(13.77))
UPPER_FOUR = (Cm(18.75), Cm(3.15))
MIDDLE_FOUR = (Cm(18.75), Cm(8.46))
BOTTOM_FOUR = (Cm(18.75), Cm(13.77))

counter = 0

for idx, folder in enumerate(sorted(all_paths)):
    if counter > 0:
        if counter % slides_per_pptx == 0:
            count_suffix = str(int(counter / slides_per_pptx))
            prs.save(home + '/spectrum_visualization_all_results' + folder_in_use[1:] + '_'+ count_suffix + '.pptx')
            prs = Presentation()
    pic_title = os.path.basename(os.path.abspath(folder))
    print(pic_title)
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = pic_title
    slide.shapes.title.font = Pt(12)
    slide_label = slide.shapes.add_textbox(Cm(5), Cm(0),1,1)
    slide_label.text = all_labels[idx]

    if glob.glob(folder + folder_in_use + '/MASK_COORDINATES.csv'):
        crop_coords(folder)


        slide.shapes.add_picture(glob.glob(folder + '/_crops/*RGB_cropped_with_mask.png')[0], MIDDLE_TWO[0], MIDDLE_TWO[1], width, height)
        slide.shapes.add_picture(glob.glob(folder + '/_crops/*Oxygenation_cropped_with_mask.png')[0], BOTTOM_TWO[0], BOTTOM_TWO[1], width, height)
    else:
        slide.shapes.add_picture(glob.glob(folder + '/_crops/*RGB-Image.png')[0], MIDDLE_TWO[0],MIDDLE_TWO[1], width, height)
        slide.shapes.add_picture(glob.glob(folder + '/_crops/*Oxygenation.png')[0], BOTTOM_TWO[0],BOTTOM_TWO[1], width, height)

    for spectrum_path in sorted(glob.glob(folder + folder_in_use + '/*fromCSV*with-scale*')):
        print(spectrum_path)
        
        if not os.path.basename(spectrum_path).startswith('.'):

            if 'spectrum_fromCSV1_' in spectrum_path:
                if '0_derivative' in spectrum_path:
                    slide.shapes.add_picture(spectrum_path, UPPER_ONE[0], UPPER_ONE[1], width, height)
                elif '1_derivative' in spectrum_path:
                    slide.shapes.add_picture(spectrum_path, MIDDLE_ONE[0], MIDDLE_ONE[1], width, height)
                elif '2_derivative' in spectrum_path:
                    slide.shapes.add_picture(spectrum_path, BOTTOM_ONE[0], BOTTOM_ONE[1], width, height)
                else:
                    slide.shapes.add_picture(spectrum_path, UPPER_ONE[0], UPPER_ONE[1], width, height)


            elif 'histogram_fromCSV1_' in spectrum_path:
                if '0_derivative' in spectrum_path:
                    slide.shapes.add_picture(spectrum_path, UPPER_TWO[0], UPPER_TWO[1], width, height)
                else:
                    slide.shapes.add_picture(spectrum_path, UPPER_TWO[0], UPPER_TWO[1], width, height)

            elif 'spectrum_fromCSV5_' in spectrum_path:
                if '0_derivative' in spectrum_path:
                    slide.shapes.add_picture(spectrum_path, UPPER_THREE[0], UPPER_THREE[1], width, height)
                elif '1_derivative' in spectrum_path:
                    slide.shapes.add_picture(spectrum_path, MIDDLE_THREE[0], MIDDLE_THREE[1], width, height)
                elif '2_derivative' in spectrum_path:
                    slide.shapes.add_picture(spectrum_path, BOTTOM_THREE[0], BOTTOM_THREE[1], width, height)
                else:
                    slide.shapes.add_picture(spectrum_path, UPPER_THREE[0], UPPER_THREE[1], width, height)

            elif 'histogram_fromCSV5_' in spectrum_path:
                if '0_derivative' in spectrum_path:
                    slide.shapes.add_picture(spectrum_path, UPPER_FOUR[0], UPPER_FOUR[1], width, height)
                elif '1_derivative' in spectrum_path:
                    slide.shapes.add_picture(spectrum_path, MIDDLE_FOUR[0], MIDDLE_FOUR[1], width, height)
                elif '2_derivative' in spectrum_path:
                    slide.shapes.add_picture(spectrum_path, BOTTOM_FOUR[0], BOTTOM_FOUR[1], width, height)
                else:
                    slide.shapes.add_picture(spectrum_path, UPPER_FOUR[0], UPPER_FOUR[1], width, height)
    counter = counter + 1
count_suffix = str(int(counter / slides_per_pptx)+1)
prs.save(home + '/spectrum_visualization_all_results' + folder_in_use[1:] + '_' + count_suffix  + '.pptx')













