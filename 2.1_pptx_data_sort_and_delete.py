"""
@author: Alexander Studier-Fischer, Jan Odenthal, Berkin Özdemir, University of Heidelberg
"""
###
# Copying Image Folders to respective "data" Directory based on Powerpoints for each Situation
# and subsequently creating summary Excel sheet of newly sorted data Folder


from pptx import Presentation
import shutil, os
import xlsxwriter
import glob
from pathlib import Path
from shutil import copyfile
import re

### Activate/Deactivate Delete ###
delete = True
####                           ###


# Core directories and lists for loops
home = os.getcwd().replace("\\", "/")

primary_paths = [,
                 ]                                                  ### Primary Paths



data_path = './data'
ppts = list()
title = list()

def check_completeness(path):
    for i in glob.glob(path + '/*/*'):
        if bool(re.match('[\d/-_]+$', os.path.basename(i))) == True:
            current_img_path = os.path.abspath(i)
            current_img = os.path.basename(i)
            print(current_img)

            for l in primary_paths:
                if not len(glob.glob(l + '/**' + '/' + current_img)) == 0:
                    primary_folder = glob.glob(l + '/**' + '/' + current_img)[0]

            for j in os.listdir(primary_folder):
                if not os.path.exists(current_img_path + '/' + j):
                    print('File not in data folder, therefore copied now.')
                    print(j)
                    copyfile(primary_folder + '/' + j, current_img_path + '/' + j)


check_completeness(data_path)

def look_for_folder(path_to_op):
    for file_to_be_found in sorted(Path(path_to_op).rglob(a)):
        file_to_be_found = os.path.abspath(file_to_be_found)
        file_to_be_found_base = os.path.basename(file_to_be_found)
        print(file_to_be_found_base)

        op_to_found = Path(file_to_be_found).parent.absolute()
        op_to_be_found_base = os.path.basename(op_to_found)

        if os.path.exists(home + '/data/' + op_to_be_found_base + '/' + file_to_be_found_base) != True:
            shutil.copytree(file_to_be_found, home + '/data/' + op_to_be_found_base + '/' + a)
            print('Copied ' + file_to_be_found_base + ' successfully!')

def check_data_parent(path):
    parent = os.path.abspath(Path(path).parent)
    if len(glob.glob(parent + '/*')) == 0:
        shutil.rmtree(parent)



def check_data(titles):
    for file in glob.glob(data_path + '/*/2*'):
        if bool(re.match('[\d/-_]+$', os.path.basename(file))) == True:
            pic_folder = os.path.basename(file)
            if not pic_folder in titles:
                print(pic_folder)
                print('File has been deleted from presentation and is now being removed from data folder')
                shutil.rmtree(file)
                check_data_parent(file)


# Find all Powerpoints in current Catalogization folder
for file in os.listdir(home):
    if "_data_sort_" in file and file.endswith(".pptx") and file.startswith('.') == False and file.startswith('~$') == False:
        print(file)
        ppts.append(os.path.join(home, file).replace("\\", "/"))

# Get contents of Powerpoints
for ppt_current in ppts:
    # Get Powerpoint
    ppt = Presentation(ppt_current)
    slides = ppt.slides

    # Get title
    for slide in slides:
        # Accessing table
        title_current = slide.shapes.title.text
        title.append(title_current)
        #print(title_current)

if delete == True:
    check_data(title)

# Look for data and copy
for a in title:
    if glob.glob(home + '/data' + '/**/*' + a):
        print('File already in Data Folder')

    else:
        print('Looking for file')

        for i in primary_paths:
            if glob.glob(i + '/**/*' + a):
                look_for_folder(i)
                current_path = i



print('All Files copied successfully!')

#Set up for Excel
data_folders = os.listdir(data_path)
row = 0
col = 0
# Create and format Excel sheet
workbook = xlsxwriter.Workbook('pptx_data_sort_results.xlsx')
bold = workbook.add_format({'bold': True})
worksheet = workbook.add_worksheet()
worksheet.set_column(0, 1, 1)
worksheet.set_column(2, 2, 40)
worksheet.set_column(3, 3, 15)
worksheet.set_column(4, 4, 1)
worksheet.set_column(5, 5, 30)
worksheet.write(row + 2, col + 2, 'Experiments', bold)
worksheet.write(row + 3, col + 2, '=COUNTA(C5:C5042)', bold)
worksheet.write(row + 2, col + 3, 'Animal Numbers', bold)
worksheet.write(row + 3, col + 3, '=COUNTA(D5:D5042)', bold)
worksheet.write(row + 2, col + 5, 'Sample Numbers', bold)
worksheet.write(row + 3, col + 5, '=COUNTA(F5:F5042)', bold)

# loop trough data folders
for folder_current in sorted(data_folders):
    if folder_current.startswith('.') == False:
        worksheet.write_string(row + 6, col + 2, folder_current)
        worksheet.write_string(row + 6, col + 3, 'x')
        temp = os.listdir(data_path + '/' + folder_current)

        for subfolder_current in sorted(temp):
            if subfolder_current.startswith('20') == True:
                worksheet.write_string(row + 6, col + 5, subfolder_current)
                row += 1

workbook.close()
