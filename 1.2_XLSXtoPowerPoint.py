"""
@author: Alexander Studier-Fischer, Jan Odenthal, Berkin Özdemir, University of Heidelberg
"""
from pptx import *
import pandas as pd
import os
import glob
import numpy as np
import datetime


###First Part Timeseries, second part ppt###

base = os.getcwd()

header = np.array(['Folder Name', 'Current Step; Situation', 'View (for Snapshot) or Timepoint (for Timeseries)',
                   'Duplicate (only if necessary))', 'Mode', 'Organ / Body Part','Description', 'Time in Case of Time Series','','','','','']).reshape(1,13)

file = os.path.abspath(glob.glob('./_TIVITARecordingsProtocol.xlsx')[0])
complete_file = np.asarray(pd.DataFrame(pd.read_excel(file, header= None)))
images= np.asarray(pd.DataFrame(pd.read_excel(file, header= None, usecols=[0])))
marks = np.asarray(pd.DataFrame(pd.read_excel(file, header= None, usecols=[12])))
marks_start = np.where(marks == 1)[0]
marks_finish = np.where(marks == 2)[0]

output_array = np.full([images.size, 2], 0.1)

# Intervals
for i in np.arange(0, marks_start.shape[0], 1):
    first_date = datetime.datetime.strptime(str(images[marks_start[i]])[2:-2], '%Y_%m_%d_%H_%M_%S')
    print(i, first_date)
    for j in np.arange(marks_start[i],marks_finish[i]+1, 1):
        next_date = datetime.datetime.strptime(str(images[j])[2:-2], '%Y_%m_%d_%H_%M_%S')
        print(next_date)
        print(next_date-first_date)
        diff_min_round = round(((next_date-first_date).total_seconds()/60)*2)/2
        diff_min = round((next_date-first_date).total_seconds()/60,2)
        output_array[j, 0] = diff_min
        output_array[j, 1] = diff_min_round

# Empty spaces
for idx, i in enumerate(output_array[:, 0]):
    if i == 0.1:
        output_array[idx,0]= 'NaN'

for idx, i in enumerate(output_array[:, 1]):
    if i == 0.1:
        output_array[idx,1]= 'NaN'


df_output = np.asarray(pd.DataFrame(output_array).replace(np.nan, '-', regex=True))
complete_file[:,7],complete_file[:,8] = df_output[:,1],df_output[:,0]
header_complete_file = np.append(header, complete_file, axis = 0)


pd.DataFrame(header_complete_file).to_excel('./_TIVITARecordingsProtocol_complete.xlsx', header = None, index= None)




# ============= INSTRUCTION ================
# 1. Make sure you have python-pptx to run this program. 
#    Install it thru pip if you do not.
# 2. Fill in the values in the INPUT section according to the examples.
# 3. Run this script.

# -------- Requirements for CSV ------------
# 1. The file ends in csv
# 2. Its columns are labelled
#   i.e col1 | col2 | col3
#       -----|------|------
#       1    | 2    | 3
#       -----|------|------
#       5    | 4    | 3
# 3. Its columns match the ones in the constants section
#    (See SHORT_COLUMNS, LONG_COLUMNS)
#    Edit it to fit as needed.

# ================ INPUT ===================
PATH_TO_EXCEL_SHEET = "./_TIVITARecordingsProtocol_complete.xlsx"
# For example:
# PATH_TO_EXCEL_SHEET = "./TIVITARecordingsProtocol21.csv"

PATH_TO_POWERPOINT = "./_PowerPoint.pptx"
# For example:
# PATH_TO_POWERPOINT = "./test.pptx"

# ============== CONSTANTS =================
# Based on ImageMuncher

# The index of the table in the powerpoint
SHORT_TABLE_INDEX = -2
LONG_TABLE_INDEX = -1

TITLE_INDEX = 0
KEY_COLUMN = "Folder Name"

SHORT_COLUMN_NUM = 6
SHORT_COLUMNS = ["Current Step; Situation", 
            "View (for Snapshot) or Timepoint (for Timeseries)", 
            "Duplicate (only if necessary))", 
            "Mode", 
            "Organ / Body Part",
            "Time in Case of Time Series"]

LONG_COLUMN_NUM = 1
LONG_COLUMNS = ["Description"]

# ==========================================

# Accessing the excel (csv) sheet
spreadsheet = pd.read_excel(PATH_TO_EXCEL_SHEET)



# Accessing the powerpoint
presentation = Presentation(PATH_TO_POWERPOINT)
slides = presentation.slides

for slide in slides:
    # Accessing table
    long_table = slide.shapes[LONG_TABLE_INDEX].table
    short_table = slide.shapes[SHORT_TABLE_INDEX].table
    # Accessing title 
    title = slide.shapes[TITLE_INDEX].text
    # Filter spreadsheet to only contain ones with the title
    mask = spreadsheet[KEY_COLUMN] == title
    filtered_sheet = spreadsheet[mask]

    for cell_index in range(SHORT_COLUMN_NUM):
        # Accessing cell
        cell = short_table.cell(0, cell_index)
        print(SHORT_COLUMNS[cell_index])
        column_value = str(filtered_sheet[SHORT_COLUMNS[cell_index]].iloc[0])
        cell.text = column_value
        print("Column Value: " + column_value)

    for cell_index in range(LONG_COLUMN_NUM):
        # Accessing cell
        cell = long_table.cell(0, cell_index)
        column_value = str(filtered_sheet[LONG_COLUMNS[cell_index]].iloc[0])
        cell.text = column_value
        print("Column Value: " + column_value)
        print("\n")

presentation.save(PATH_TO_POWERPOINT)

