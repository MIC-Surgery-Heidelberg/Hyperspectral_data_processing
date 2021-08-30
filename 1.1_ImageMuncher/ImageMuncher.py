"""
@author: Alexander Studier-Fischer, Isabella Camplisson, Berkin Özdemir, University of Heidelberg
"""
from pptx import *
from INPUTS import *
from utility import *
import os

def main():
    # Make a new presentation
    prs = Presentation()

    sub_folder_paths = get_sub_folder_paths(PATH_TO_DIRECTORIES)

    caption_text = get_caption_text(PATH_TO_DIRECTORIES)

    # Make a new slide for each sub folder and its contents
    for sub_folder_path in sub_folder_paths:
        sub_folder_content_paths = get_sub_folder_content_paths(sub_folder_path)

        # Make a new slide
        slide = prs.slides.add_slide(prs.slide_layouts[SLIDE_LAYOUT_TITLE_AND_CONTENT])  # title and content layout

        # Add the background image
        # add_background(prs, slide, PATH_TO_BACKGROUND_PIC)

        # Add the title - using the sub folder path
        title_text = get_date(sub_folder_path)
        print("Processing: " + str(title_text))
        add_title(slide, title_text)

        # Add the images
        for i in image_dic:
            if i in sub_folder_content_paths:
                add_picture(slide, sub_folder_content_paths[i], image_dic[i], MAIN_IMAGE_SIZE)

        # Add the caption text
        add_text(slide, caption_text, CAPTION_TXT_LOC, CAPTION_IMAGE_SIZE)

        # Add the table
        add_table(slide, SHORT_TABLE_ROW_NUM, SHORT_TABLE_COL_NUM, SHORT_TABLE_LOC, SHORT_TABLE_SIZE)
        add_table(slide, LONG_TABLE_ROW_NUM, LONG_TABLE_COL_NUM, LONG_TABLE_LOC, LONG_TABLE_SIZE)


    prs.save(PATH_TO_POWERPOINT)

main()
